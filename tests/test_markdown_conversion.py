from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import threading
import unittest
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
SCRIPT_DIR = ROOT / "skills" / "markdown-conversion" / "scripts"
PNG_1X1 = bytes.fromhex(
    "89504E470D0A1A0A0000000D4948445200000001000000010802000000907753DE"
    "0000000C49444154789C63606060000000040001F61738550000000049454E44AE426082"
)


def run_script(*args: str, cwd: Path | None = None) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, *args],
        cwd=str(cwd or ROOT),
        text=True,
        capture_output=True,
        check=True,
    )


class MarkdownConversionTests(unittest.TestCase):
    def test_html_conversion_writes_manifest_and_profile(self) -> None:
        try:
            import bs4  # noqa: F401
            import markdownify  # noqa: F401
        except ImportError as exc:
            self.skipTest(f"missing optional dependency: {exc.name}")

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            html = tmp_path / "sample.html"
            html.write_text(
                '<html><body><h1>Doc Test</h1><p>Hello</p>'
                '<img alt="dot" src="data:image/png;base64,'
                'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAIAAACQd1PeAAAADGlEQVR4nGNgYGAEAAAEAAH2Fj8VAAAAAElFTkSuQmCC'
                '"></body></html>',
                encoding="utf-8",
            )
            out = tmp_path / "out.md"

            run_script(str(SCRIPT_DIR / "doc_to_md.py"), str(html), "-o", str(out))

            self.assertIn("![dot](out_files/image_001.png)", out.read_text(encoding="utf-8"))
            manifest = json.loads((tmp_path / "out_files" / "image_manifest.json").read_text(encoding="utf-8"))
            self.assertEqual(manifest[0]["source_kind"], "html_image")
            profile = json.loads((tmp_path / "source_profile.json").read_text(encoding="utf-8"))
            self.assertEqual(profile["outputs"]["image_count"], 1)
            self.assertEqual(profile["markdown"]["image_ref_count"], 1)

    def test_ppt_conversion_exports_chart_data_and_deduplicates_images(self) -> None:
        try:
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.util import Inches
        except ImportError as exc:
            self.skipTest(f"missing optional dependency: {exc.name}")

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            png = tmp_path / "dot.png"
            png.write_bytes(PNG_1X1)

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Quarterly Metrics"
            chart_data = CategoryChartData()
            chart_data.categories = ["Q1", "Q2"]
            chart_data.add_series("Revenue", (10, 20))
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(4),
                Inches(1.3),
                Inches(4.5),
                Inches(3),
                chart_data,
            )
            slide.shapes.add_picture(str(png), Inches(0.7), Inches(3.4), Inches(0.5), Inches(0.5))
            slide.shapes.add_picture(str(png), Inches(1.4), Inches(3.4), Inches(0.5), Inches(0.5))
            deck = tmp_path / "sample.pptx"
            prs.save(deck)
            out = tmp_path / "out.md"

            run_script(str(SCRIPT_DIR / "ppt_to_md.py"), str(deck), "-o", str(out))

            markdown = out.read_text(encoding="utf-8")
            self.assertIn("| Category | Revenue |", markdown)
            self.assertIn("| Q2 | 20 |", markdown)
            manifest = json.loads((tmp_path / "out_files" / "image_manifest.json").read_text(encoding="utf-8"))
            self.assertEqual(len(manifest), 1)
            self.assertEqual(manifest[0]["usage_count"], 2)
            self.assertEqual(len(manifest[0]["occurrences"]), 2)

    def test_web_conversion_writes_image_sources(self) -> None:
        try:
            import bs4  # noqa: F401
        except ImportError as exc:
            self.skipTest(f"missing optional dependency: {exc.name}")

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            (tmp_path / "dot.png").write_bytes(PNG_1X1)
            (tmp_path / "index.html").write_text(
                '<html><body><main><h1>Web Test</h1><p>This is enough content for conversion.</p>'
                '<img alt="dot" src="/dot.png"></main></body></html>',
                encoding="utf-8",
            )

            class QuietHandler(SimpleHTTPRequestHandler):
                def log_message(self, format: str, *args: object) -> None:
                    return

            server = ThreadingHTTPServer(("127.0.0.1", 0), QuietHandler)
            server_dir = tmp_path
            old_cwd = Path.cwd()
            try:
                import os
                os.chdir(server_dir)
                thread = threading.Thread(target=server.serve_forever, daemon=True)
                thread.start()
                url = f"http://127.0.0.1:{server.server_port}/index.html"
                out = tmp_path / "web.md"
                run_script(str(SCRIPT_DIR / "web_to_md.py"), url, "-o", str(out), "--raw", cwd=ROOT)
            finally:
                server.shutdown()
                server.server_close()
                os.chdir(old_cwd)

            sources = json.loads((tmp_path / "web_files" / "image_sources.json").read_text(encoding="utf-8"))
            self.assertEqual(sources["items"][0]["license_status"], "unknown")
            self.assertIn("/dot.png", sources["items"][0]["download_url"])

    def test_convert_json_output_and_source_profile(self) -> None:
        try:
            import bs4  # noqa: F401
            import markdownify  # noqa: F401
        except ImportError as exc:
            self.skipTest(f"missing optional dependency: {exc.name}")

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            html = tmp_path / "sample.html"
            html.write_text("<html><body><h1>JSON Test</h1><p>Hello</p></body></html>", encoding="utf-8")
            out = tmp_path / "json.md"

            result = run_script(
                str(SCRIPT_DIR / "convert.py"),
                str(html),
                "-o",
                str(out),
                "--json",
            )

            payload = json.loads(result.stdout.strip().splitlines()[-1])
            self.assertEqual(payload["markdown"], str(out.resolve()))
            self.assertTrue((tmp_path / "source_profile.json").exists())


if __name__ == "__main__":
    unittest.main()
