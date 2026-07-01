#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter

Extracts slide text, hyperlinks, tables, speaker notes, and embedded pictures from
Open XML PowerPoint files into Markdown.

Primary use case: PPTX source decks -> Markdown for PPT generation input.

Dependency:
    pip install python-pptx
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from urllib.parse import quote

from pptx import Presentation
from pptx.enum.action import PP_ACTION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

sys.path.insert(0, str(Path(__file__).parent))
from _conversion_profile import write_source_profile  # noqa: E402
from _image_filter import should_keep_image as _should_keep_image  # noqa: E402
from _image_filter import should_keep_image_bytes as _should_keep_image_bytes  # noqa: E402


EMU_PER_INCH = 914400
UNSUPPORTED_URL_SCHEMES = ("javascript:", "data:", "vbscript:", "file:")
SUPPORTED_FORMATS = {
    ".pptx": "PowerPoint Presentation",
    ".pptm": "Macro-enabled PowerPoint Presentation",
    ".ppsx": "PowerPoint Slide Show",
    ".ppsm": "Macro-enabled PowerPoint Slide Show",
    ".potx": "PowerPoint Template",
    ".potm": "Macro-enabled PowerPoint Template",
}


@dataclass
class LeafShape:
    """Flattened leaf shape with stable position ordering."""

    shape: object
    top: int
    left: int


@dataclass
class SavedImage:
    """Extracted image asset plus manifest metadata."""

    filename: str
    manifest_entry: dict[str, object]
    is_new_asset: bool


def normalize_text(value: str) -> str:
    """Collapse whitespace while preserving paragraph boundaries elsewhere."""
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"\s+", " ", line).strip() for line in value.split("\n")]
    lines = [line for line in lines if line]
    return "\n".join(lines)


def normalize_inline_text(value: str) -> str:
    """Collapse run-level whitespace while preserving leading/trailing spaces."""
    value = value.replace("\r\n", " ").replace("\r", " ").replace("\n", " ")
    return re.sub(r"\s+", " ", value)


def escape_markdown_link_text(value: str) -> str:
    """Escape Markdown link label delimiters."""
    return value.replace("\\", "\\\\").replace("[", r"\[").replace("]", r"\]")


def escape_markdown_link_target(value: str) -> str:
    """Escape Markdown inline-link target delimiters."""
    return quote(value.strip(), safe="/:?=&%#@!$'*+,;")


def markdown_link(label: str, target: str) -> str:
    """Return a Markdown inline link."""
    return f"[{escape_markdown_link_text(label)}]({escape_markdown_link_target(target)})"


def escape_table_cell(value: str) -> str:
    """Escape Markdown table syntax inside a cell."""
    return normalize_text(value).replace("|", r"\|") or " "


def safe_position(shape: object, attr: str) -> int:
    """Read shape position safely, tolerating broken placeholder inheritance."""
    try:
        return int(getattr(shape, attr, 0) or 0)
    except Exception:
        return 0


def supported_url(url: str | None) -> str | None:
    """Return a safe URL for Markdown output, or None."""
    if not url:
        return None
    value = str(url).strip()
    if not value:
        return None
    if any(value.lower().startswith(scheme) for scheme in UNSUPPORTED_URL_SCHEMES):
        return None
    return value


def iter_leaf_shapes(shapes: object) -> list[LeafShape]:
    """Return a flattened, reading-order list of shapes."""
    items: list[LeafShape] = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            items.extend(iter_leaf_shapes(shape.shapes))
            continue
        items.append(
            LeafShape(
                shape=shape,
                top=safe_position(shape, "top"),
                left=safe_position(shape, "left"),
            )
        )
    items.sort(key=lambda item: (item.top, item.left))
    return items


def hyperlink_address(item: object) -> str | None:
    """Return an external hyperlink address for a run or shape, if present."""
    try:
        if hasattr(item, "click_action"):
            action = item.click_action
            if action.action == PP_ACTION.NAMED_SLIDE:
                target_slide = action.target_slide
                if target_slide is not None:
                    prs = item.part.slide.part.package.presentation_part.presentation
                    return f"#slide-{list(prs.slides).index(target_slide) + 1}"
            address = action.hyperlink.address
        else:
            address = item.hyperlink.address
    except Exception:
        return None
    return supported_url(address)


def resolve_run_internal_jump(run: object, shape: object | None) -> str | None:
    """Return #slide-N for a run-level slide jump, if present."""
    if shape is None:
        return None
    r_id = ""
    try:
        rpr = run._r.find(qn("a:rPr"))
        if rpr is None:
            return None
        hlink = rpr.find(qn("a:hlinkClick"))
        if hlink is None or "hlinksldjump" not in (hlink.get("action", "") or ""):
            return None
        r_id = hlink.get(qn("r:id"), "")
        if not r_id:
            return None
        target_slide = shape.part.related_part(r_id).slide
        prs = shape.part.slide.part.package.presentation_part.presentation
        return f"#slide-{list(prs.slides).index(target_slide) + 1}"
    except Exception:
        print(f"[WARN] ppt_to_md: could not resolve slide jump rId={r_id}", file=sys.stderr)
        return None


def run_hyperlink(run: object, shape: object | None) -> str | None:
    """Return a run's external or internal hyperlink target."""
    return resolve_run_internal_jump(run, shape) or hyperlink_address(run)


def paragraph_has_hyperlink(paragraph: object) -> bool:
    """Return whether a paragraph contains any run-level hyperlink."""
    for run in getattr(paragraph, "runs", []):
        try:
            if run.hyperlink.address:
                return True
        except AttributeError:
            pass
        try:
            rpr = run._r.find(qn("a:rPr"))
            if rpr is not None and rpr.find(qn("a:hlinkClick")) is not None:
                return True
        except AttributeError:
            continue
    return False


def paragraph_to_markdown(paragraph: object, fallback_hyperlink: str | None = None, shape: object | None = None) -> str:
    """Convert one PowerPoint paragraph into Markdown text."""
    parts: list[str] = []
    current_url: str | None = None
    current_text = ""
    has_run_hyperlink = False

    def flush() -> None:
        nonlocal current_text
        if not current_text:
            return
        if current_url is None:
            parts.append(current_text)
            current_text = ""
            return
        leading = current_text[: len(current_text) - len(current_text.lstrip())]
        trailing = current_text[len(current_text.rstrip()):]
        core = current_text.strip()
        display = core if core else current_url
        parts.append(f"{leading}{markdown_link(display, current_url)}{trailing}")
        current_text = ""

    for run in getattr(paragraph, "runs", []):
        text = normalize_inline_text(getattr(run, "text", ""))
        if not text:
            continue

        address = run_hyperlink(run, shape)
        if address:
            has_run_hyperlink = True
        if address != current_url:
            flush()
            current_url = address
        current_text += text
    flush()

    text_md = re.sub(r"\s+", " ", "".join(parts)).strip()
    if not text_md:
        text_md = normalize_text(getattr(paragraph, "text", "")).replace("\n", " ")

    if fallback_hyperlink and text_md and not has_run_hyperlink:
        return markdown_link(text_md, fallback_hyperlink)
    return text_md


def text_frame_to_markdown(text_frame: object, fallback_hyperlink: str | None = None, shape: object | None = None) -> str:
    """Convert a PowerPoint text frame into Markdown."""
    paragraphs = []
    visible_paragraphs = []
    for paragraph in text_frame.paragraphs:
        text_md = paragraph_to_markdown(paragraph, fallback_hyperlink, shape)
        if text_md or paragraph_has_hyperlink(paragraph):
            visible_paragraphs.append((paragraph, text_md))

    if not visible_paragraphs:
        return ""

    list_like = any(paragraph.level > 0 for paragraph, _ in visible_paragraphs)
    if not list_like:
        list_like = len(visible_paragraphs) > 1

    for paragraph, text in visible_paragraphs:
        if list_like:
            indent = "  " * max(paragraph.level, 0)
            paragraphs.append(f"{indent}- {text}")
        else:
            paragraphs.append(text)

    if list_like:
        return "\n".join(paragraphs)
    return "\n\n".join(paragraphs)


def table_to_markdown(table: object) -> str:
    """Convert a PowerPoint table to a Markdown table."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            try:
                cell_text = text_frame_to_markdown(cell.text_frame)
            except Exception:
                cell_text = getattr(cell, "text", "")
            cells.append(escape_table_cell(cell_text.replace("\n", " ")))
        rows.append(cells)

    if not rows:
        return ""

    column_count = max(len(row) for row in rows)
    normalized_rows = [row + [" "] * (column_count - len(row)) for row in rows]
    header = normalized_rows[0]
    separator = ["---"] * column_count
    body = normalized_rows[1:]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def format_chart_value(value: object) -> str:
    """Render a chart data point, trimming whole-number floats."""
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def chart_to_markdown(chart: object, name: str) -> str:
    """Convert readable native PowerPoint chart data to Markdown."""
    try:
        chart_type = str(chart.chart_type)
    except (ValueError, AttributeError):
        chart_type = ""

    categories: list[str] = []
    try:
        plots = list(chart.plots)
        if plots:
            categories = [
                escape_table_cell(str(cat)) if cat is not None else ""
                for cat in plots[0].categories
            ]
    except (ValueError, IndexError, AttributeError):
        categories = []

    series_data: list[tuple[str, list[object]]] = []
    try:
        for index, series in enumerate(chart.series, start=1):
            try:
                values = list(series.values)
            except (ValueError, TypeError, AttributeError):
                values = []
            label = str(series.name) if getattr(series, "name", None) else f"Series {index}"
            series_data.append((escape_table_cell(label), values))
    except (ValueError, AttributeError):
        series_data = []

    header = f"> [Chart] {name}" + (f" - {chart_type}" if chart_type else "")
    row_count = len(categories) if categories else max((len(v) for _, v in series_data), default=0)
    if not series_data or row_count == 0:
        return header

    table_header = (["Category"] if categories else ["#"]) + [series_name for series_name, _ in series_data]
    lines = [
        header,
        "",
        "| " + " | ".join(table_header) + " |",
        "| " + " | ".join(["---"] * len(table_header)) + " |",
    ]
    for row_index in range(row_count):
        if categories:
            label = categories[row_index] if row_index < len(categories) else ""
        else:
            label = str(row_index + 1)
        cells = [label]
        for _, values in series_data:
            cells.append(format_chart_value(values[row_index]) if row_index < len(values) else "")
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def image_size_from_bytes(blob: bytes) -> tuple[int | None, int | None]:
    """Return bitmap dimensions when Pillow can decode the bytes."""
    try:
        from PIL import Image
        with Image.open(BytesIO(blob)) as image:
            return image.width, image.height
    except Exception:
        return None, None


def image_manifest_entry(
    shape: object,
    image: object,
    filename: str,
    slide_index: int,
    image_index: int,
) -> dict[str, object]:
    """Build metadata for one extracted PowerPoint image."""
    blob = bytes(getattr(image, "blob", b""))
    pixel_width, pixel_height = image_size_from_bytes(blob)
    pixel_ratio = (
        pixel_width / pixel_height
        if pixel_width and pixel_height
        else None
    )
    display_width = safe_position(shape, "width")
    display_height = safe_position(shape, "height")
    display_ratio = (
        display_width / display_height
        if display_width > 0 and display_height > 0
        else None
    )
    ext = (getattr(image, "ext", "") or Path(filename).suffix.lstrip(".") or "bin").lower()
    asset_kind = "office_vector" if ext in {"emf", "wmf"} else "bitmap"
    return {
        "index": image_index,
        "filename": filename,
        "asset_kind": asset_kind,
        "svg_renderable": asset_kind != "office_vector",
        "pptx_native_supported": True,
        "source_kind": "pptx_picture",
        "source_ext": f".{ext}",
        "content_type": str(getattr(image, "content_type", "")),
        "pixel_width": pixel_width,
        "pixel_height": pixel_height,
        "pixel_ratio": round(pixel_ratio, 6) if pixel_ratio else None,
        "usage_count": 1,
        "occurrences": [
            {
                "slide_index": slide_index,
                "shape_name": str(getattr(shape, "name", "")),
                "display_left_emu": safe_position(shape, "left"),
                "display_top_emu": safe_position(shape, "top"),
                "display_width_emu": display_width,
                "display_height_emu": display_height,
                "display_width_in": round(display_width / EMU_PER_INCH, 4) if display_width else None,
                "display_height_in": round(display_height / EMU_PER_INCH, 4) if display_height else None,
                "display_ratio": round(display_ratio, 6) if display_ratio else None,
            }
        ],
    }


def image_occurrence(shape: object, slide_index: int) -> dict[str, object]:
    """Build one image occurrence row for the manifest."""
    display_width = safe_position(shape, "width")
    display_height = safe_position(shape, "height")
    display_ratio = (
        display_width / display_height
        if display_width > 0 and display_height > 0
        else None
    )
    return {
        "slide_index": slide_index,
        "shape_name": str(getattr(shape, "name", "")),
        "display_left_emu": safe_position(shape, "left"),
        "display_top_emu": safe_position(shape, "top"),
        "display_width_emu": display_width,
        "display_height_emu": display_height,
        "display_width_in": round(display_width / EMU_PER_INCH, 4) if display_width else None,
        "display_height_in": round(display_height / EMU_PER_INCH, 4) if display_height else None,
        "display_ratio": round(display_ratio, 6) if display_ratio else None,
    }


def add_image_occurrence(entry: dict[str, object], occurrence: dict[str, object]) -> None:
    """Append an image occurrence and refresh aggregate manifest fields."""
    occurrences = entry.setdefault("occurrences", [])
    if isinstance(occurrences, list):
        occurrences.append(occurrence)
        entry["usage_count"] = len(occurrences)


def image_cache_key(image: object) -> str:
    """Return a stable content hash for deduplicating image assets."""
    return hashlib.sha256(bytes(getattr(image, "blob", b""))).hexdigest()


def save_picture(
    shape: object,
    asset_dir: Path,
    slide_index: int,
    image_index: int,
    asset_cache: dict[str, SavedImage],
) -> SavedImage | None:
    """Persist a picture shape to the output asset directory."""
    try:
        image = shape.image
    except Exception:
        return None

    cache_key = image_cache_key(image)
    cached = asset_cache.get(cache_key)
    if cached is not None:
        add_image_occurrence(cached.manifest_entry, image_occurrence(shape, slide_index))
        return SavedImage(
            filename=cached.filename,
            manifest_entry=cached.manifest_entry,
            is_new_asset=False,
        )

    ext = (image.ext or "png").lower()
    filename = f"slide_{slide_index:02d}_image_{image_index:02d}.{ext}"
    output_path = asset_dir / filename
    output_path.write_bytes(image.blob)
    saved = SavedImage(
        filename=filename,
        manifest_entry=image_manifest_entry(shape, image, filename, slide_index, image_index),
        is_new_asset=True,
    )
    asset_cache[cache_key] = saved
    return saved


def _shape_passes_ai_filter(
    shape: object,
    slide_size: tuple[int, int],
    seen_hashes: set[str],
) -> bool:
    """Run the shared decorative-image filter against a PPT picture shape."""
    try:
        image = shape.image
    except Exception:
        return True  # No accessible image (e.g. linked picture) — keep ref
    blob = image.blob
    try:
        from PIL import Image
        import io
        with Image.open(io.BytesIO(blob)) as img:
            pixel_w, pixel_h = img.size
    except Exception:
        return _should_keep_image_bytes(blob, seen_hashes=seen_hashes)
    render_w = int(getattr(shape, "width", 0) or 0)
    render_h = int(getattr(shape, "height", 0) or 0)
    return _should_keep_image(
        blob,
        pixel_w or 1,
        pixel_h or 1,
        page_size=slide_size if all(slide_size) else None,
        render_size=(render_w, render_h) if (render_w and render_h) else None,
        seen_hashes=seen_hashes,
    )


def extract_notes(slide: object) -> str:
    """Extract speaker notes text from a slide, if available."""
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""

    blocks = []
    for item in iter_leaf_shapes(notes_slide.shapes):
        shape = item.shape
        if not getattr(shape, "has_text_frame", False):
            continue
        text = text_frame_to_markdown(shape.text_frame, shape=shape)
        if text:
            blocks.append(text)

    return "\n\n".join(blocks).strip()


def convert_presentation_to_markdown(
    input_path: str,
    output_path: str | None = None,
    no_images: bool = False,
    filter_images: bool = False,
) -> str:
    """Convert a supported PowerPoint file to Markdown."""
    input_file = Path(input_path)
    if not input_file.exists():
        print(f"[ERROR] File not found: {input_path}")
        return ""

    suffix = input_file.suffix.lower()
    if suffix not in SUPPORTED_FORMATS:
        supported = ", ".join(sorted(SUPPORTED_FORMATS.keys()))
        print(f"[ERROR] Unsupported format: {suffix}")
        print(f"   Supported: {supported}")
        print("   Legacy .ppt files should be resaved as .pptx or exported to PDF first.")
        return ""

    print(f"[INFO] Converting {SUPPORTED_FORMATS[suffix]}: {input_file.name}")

    if output_path:
        out_file = Path(output_path)
    else:
        out_file = input_file.with_suffix(".md")

    out_file.parent.mkdir(parents=True, exist_ok=True)
    asset_dir = out_file.parent / f"{out_file.stem}_files"

    presentation = Presentation(str(input_file))
    lines = [
        f"# {input_file.stem}",
        "",
        f"- Source: `{input_file.name}`",
        f"- Total slides: {len(presentation.slides)}",
        "",
    ]

    image_count = 0
    image_ref_count = 0
    asset_dir_used = False
    image_manifest: list[dict[str, object]] = []
    asset_cache: dict[str, SavedImage] = {}
    slide_size = (
        int(getattr(presentation, "slide_width", 0) or 0),
        int(getattr(presentation, "slide_height", 0) or 0),
    )
    seen_image_hashes: set[str] = set()

    for slide_index, slide in enumerate(presentation.slides, 1):
        lines.append(f"## Slide {slide_index}")
        lines.append("")

        blocks = []
        for item in iter_leaf_shapes(slide.shapes):
            shape = item.shape
            shape_link = hyperlink_address(shape)

            if getattr(shape, "has_table", False):
                table_md = table_to_markdown(shape.table)
                if table_md:
                    blocks.append(table_md)
                if shape_link:
                    blocks.append(f"> {markdown_link(getattr(shape, 'name', 'Table'), shape_link)}")
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if no_images:
                    if shape_link:
                        blocks.append(f"> {markdown_link(getattr(shape, 'name', 'Linked image'), shape_link)}")
                    continue
                if filter_images and not _shape_passes_ai_filter(shape, slide_size, seen_image_hashes):
                    if shape_link:
                        blocks.append(f"> {markdown_link(getattr(shape, 'name', 'Filtered linked image'), shape_link)}")
                    continue
                next_image_index = image_count + 1
                asset_dir.mkdir(parents=True, exist_ok=True)
                saved_image = save_picture(shape, asset_dir, slide_index, next_image_index, asset_cache)
                if saved_image is None:
                    label = f"Image: {getattr(shape, 'name', 'Picture')}"
                    blocks.append(
                        f"> {markdown_link(label, shape_link)}"
                        if shape_link
                        else f"> [Image] {getattr(shape, 'name', 'Picture')}"
                    )
                    continue

                image_ref_count += 1
                if saved_image.is_new_asset:
                    image_count = next_image_index
                    image_manifest.append(saved_image.manifest_entry)
                asset_dir_used = True
                image_md = f"![Slide {slide_index} Image {image_ref_count}]({asset_dir.name}/{saved_image.filename})"
                blocks.append(f"[{image_md}]({escape_markdown_link_target(shape_link)})" if shape_link else image_md)
                continue

            if getattr(shape, "has_text_frame", False):
                text_md = text_frame_to_markdown(shape.text_frame, fallback_hyperlink=shape_link, shape=shape)
                if text_md:
                    blocks.append(text_md)
                continue

            if getattr(shape, "has_chart", False):
                chart_name = getattr(shape, "name", "Chart")
                try:
                    chart_md = chart_to_markdown(shape.chart, chart_name)
                except (ValueError, AttributeError, KeyError):
                    label = f"Chart: {chart_name}"
                    chart_md = (
                        f"> {markdown_link(label, shape_link)}"
                        if shape_link
                        else f"> [Chart] {chart_name}"
                    )
                blocks.append(chart_md)
                continue

            if shape_link:
                blocks.append(f"> {markdown_link(getattr(shape, 'name', 'Linked shape'), shape_link)}")

        if blocks:
            lines.append("\n\n".join(blocks))
            lines.append("")
        else:
            lines.append("_No extractable text content._")
            lines.append("")

        notes_md = extract_notes(slide)
        if notes_md:
            lines.append("### Speaker Notes")
            lines.append("")
            lines.append(notes_md)
            lines.append("")

    markdown_content = "\n".join(lines).strip() + "\n"
    out_file.write_text(markdown_content, encoding="utf-8")
    if image_manifest:
        (asset_dir / "image_manifest.json").write_text(
            json.dumps(image_manifest, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
    profile_path = write_source_profile(
        input_path=str(input_file),
        markdown_path=str(out_file),
        converter="ppt_to_md.py",
        conversion_type="pptx",
    )

    print(f"[OK] Saved Markdown to: {out_file}")
    if asset_dir_used:
        media_files = [
            path for path in asset_dir.iterdir()
            if path.is_file() and path.name != "image_manifest.json"
        ]
        print(f"   Extracted {len(media_files)} image file(s) -> {asset_dir}")
        if image_ref_count != len(media_files):
            print(
                f"   Deduplicated {image_ref_count} image reference(s) "
                f"into {len(media_files)} asset file(s)"
            )
        if image_manifest:
            print(f"   Wrote image manifest -> {asset_dir / 'image_manifest.json'}")
    print(f"   Wrote conversion profile -> {profile_path}")

    return markdown_content


def process_directory(input_dir: str, output_dir: str | None = None, no_images: bool = False, filter_images: bool = False) -> None:
    """Convert all supported PowerPoint files in a directory to Markdown."""
    input_path = Path(input_dir)

    if output_dir:
        output_path = Path(output_dir)
    else:
        output_path = input_path

    presentation_files = sorted(
        path for path in input_path.iterdir()
        if path.is_file() and path.suffix.lower() in SUPPORTED_FORMATS
    )

    print(f"Found {len(presentation_files)} PowerPoint files")

    for presentation_file in presentation_files:
        output_file = output_path / f"{presentation_file.stem}.md"
        print(f"Processing: {presentation_file.name}")
        result = convert_presentation_to_markdown(str(presentation_file), str(output_file), no_images=no_images, filter_images=filter_images)
        if not result:
            print(f"[WARN] Skipped failed file: {presentation_file.name}")


def main() -> None:
    """Run the CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint files to Markdown",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python3 ppt_to_md.py slides.pptx
  python3 ppt_to_md.py slides.pptx -o output.md
  python3 ppt_to_md.py ./decks
  python3 ppt_to_md.py ./decks -o ./markdown
  python3 ppt_to_md.py deck.ppsx -o notes/deck.md

Supported formats:
  .pptx  .pptm  .ppsx  .ppsm  .potx  .potm

Legacy .ppt is not parsed directly. Resave it as .pptx or export it to PDF first.
        """,
    )
    parser.add_argument("input", help="Input PowerPoint file or directory")
    parser.add_argument("-o", "--output", help="Output Markdown file or directory path")
    parser.add_argument(
        "--no-images",
        action="store_true",
        help="Skip picture shapes; no asset directory or image references",
    )
    parser.add_argument(
        "--filter-images",
        action="store_true",
        help="Filter decorative images (master backgrounds, logos, low-info blocks)",
    )
    parser.add_argument(
        "--raw",
        action="store_true",
        help="Faithful reproduction (no-op for pptx; reserved for parity)",
    )

    args = parser.parse_args()

    if args.no_images and args.filter_images:
        print("Error: --no-images and --filter-images are mutually exclusive.")
        sys.exit(2)

    input_path = Path(args.input)

    if input_path.is_file():
        output = args.output or str(input_path.with_suffix(".md"))
        result = convert_presentation_to_markdown(str(input_path), output, no_images=args.no_images, filter_images=args.filter_images)
        sys.exit(0 if result else 1)
    if input_path.is_dir():
        process_directory(str(input_path), args.output, no_images=args.no_images, filter_images=args.filter_images)
        sys.exit(0)

    print(f"Error: File or directory not found: {args.input}")
    sys.exit(1)


if __name__ == "__main__":
    main()
